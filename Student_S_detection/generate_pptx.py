"""
Generate a professional project presentation (12-15 slides) for the
Student Stress Detection System using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def setup_presentation():
    """Create presentation with custom design"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title, subtitle, date=None):
    """Add professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # Dark blue-gray
    
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(52, 73, 94)
    
    if date:
        # Add date at bottom
        left = Inches(5)
        top = Inches(6.5)
        width = Inches(4)
        height = Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = date
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content_items, is_bullet=True):
    """Add content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Content
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    for idx, item in enumerate(content_items):
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16) if len(content_items) <= 4 else Pt(14)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_after = Pt(8)
        if idx == 0 and is_bullet:
            p.font.bold = True


def add_two_column_slide(prs, title, left_items, right_items):
    """Add two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Left column
    left_shape = slide.shapes.placeholders[1]
    left_frame = left_shape.text_frame
    left_frame.clear()
    left_frame.word_wrap = True
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_after = Pt(6)
    
    # Right column
    right_shape = slide.shapes.placeholders[2]
    right_frame = right_shape.text_frame
    right_frame.clear()
    right_frame.word_wrap = True
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_after = Pt(6)


def add_metrics_slide(prs, title, metrics):
    """Add slide with key metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Metrics
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for name, value in metrics:
        p = text_frame.add_paragraph()
        p.text = f"{name}: {value}"
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_after = Pt(12)
        if "Selected" in name or "Ensemble" in name:
            p.font.bold = True
            p.font.color.rgb = RGBColor(39, 174, 96)  # Green for selected model


def add_methodology_slide(prs, title, steps):
    """Add methodology/process slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Steps
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for idx, step in enumerate(steps, 1):
        p = text_frame.add_paragraph()
        p.text = f"{idx}. {step}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_after = Pt(10)
        p.font.bold = True if idx == 1 else False


def add_feature_table_slide(prs, title, features):
    """Add slide with feature table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Features in two columns
    col1_left = Inches(0.8)
    col2_left = Inches(5.2)
    start_top = Inches(1.8)
    line_height = Inches(0.4)
    
    for idx, feature in enumerate(features):
        col = col1_left if idx < len(features) // 2 else col2_left
        row = idx if idx < len(features) // 2 else idx - len(features) // 2
        
        textbox = slide.shapes.add_textbox(col, start_top + row * line_height, Inches(4), line_height)
        text_frame = textbox.text_frame
        text_frame.text = f"• {feature}"
        text_frame.paragraphs[0].font.size = Pt(13)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)


def add_conclusion_slide(prs, title, key_points, final_message):
    """Add conclusion slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
    
    # Key points
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for point in key_points:
        p = text_frame.add_paragraph()
        p.text = f"✓ {point}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(39, 174, 96)  # Green checkmarks
        p.space_after = Pt(10)
    
    # Final message
    if final_message:
        p = text_frame.add_paragraph()
        p.text = final_message
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(20)


def build_presentation(output_path="PROJECT_DECK.pptx"):
    """Build the complete presentation"""
    prs = setup_presentation()
    
    from datetime import datetime
    date_str = datetime.now().strftime("%B %Y")
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Student Stress Detection System",
        "Using Ensemble Machine Learning",
        date_str
    )
    
    # Slide 2: Problem Statement
    add_content_slide(
        prs,
        "Problem Statement",
        [
            "Student stress is a critical concern affecting academic performance and well-being",
            "Traditional assessment methods are time-consuming and subjective",
            "Need for early detection and proactive intervention",
            "Machine learning offers automated, data-driven stress prediction"
        ]
    )
    
    # Slide 3: Objectives
    add_content_slide(
        prs,
        "Project Objectives",
        [
            "Develop accurate ML system for stress level classification",
            "Compare multiple algorithms to identify best approach",
            "Implement ensemble learning for improved robustness",
            "Create user-friendly web application for real-time predictions",
            "Provide actionable recommendations for stress management"
        ]
    )
    
    # Slide 4: Dataset Overview
    add_content_slide(
        prs,
        "Dataset Overview",
        [
            "Synthetic dataset: 10,000 student records",
            "10 features: study hours, sleep, exercise, social activities, deadlines, exam pressure, family support, financial stress, academic performance, workload",
            "Target variable: Stress level (Low, Medium, High)",
            "Train-Test Split: 80-20 with stratification",
            "Balanced using SMOTE to handle class imbalance"
        ]
    )
    
    # Slide 5: Key Features
    add_feature_table_slide(
        prs,
        "Key Features Analyzed",
        [
            "Study Hours (0-12 hrs/day)",
            "Sleep Hours (3-10 hrs/night)",
            "Exercise Hours (0-10 hrs/week)",
            "Social Activities (0-7/week)",
            "Assignment Deadlines (0-10)",
            "Exam Pressure (1-10 scale)",
            "Family Support (1-10 scale)",
            "Financial Stress (1-10 scale)",
            "Academic Performance (0-100)",
            "Workload Level (1-10 scale)"
        ]
    )
    
    # Slide 6: Preprocessing Pipeline
    add_methodology_slide(
        prs,
        "Data Preprocessing Pipeline",
        [
            "Data Cleaning: Remove missing values, duplicates, and outliers (IQR method)",
            "Data Balancing: Apply SMOTE to address class imbalance",
            "Feature Scaling: StandardScaler for normalization",
            "Label Encoding: Convert categorical target to numerical",
            "Train-Test Split: 80-20 stratified split"
        ]
    )
    
    # Slide 7: Models Chosen
    add_two_column_slide(
        prs,
        "Machine Learning Models",
        [
            "Logistic Regression - Interpretable baseline model",
            "Support Vector Machine (SVM) - Non-linear classification",
            "Naive Bayes - Fast probabilistic classifier",
            "Random Forest - Captures complex feature interactions"
        ],
        [
            "Ensemble (Voting Classifier) - Combines all models",
            "Soft Voting - Uses probability predictions",
            "Improved robustness and generalization",
            "Reduced overfitting risk"
        ]
    )
    
    # Slide 8: Why Ensemble Model
    add_content_slide(
        prs,
        "Why Ensemble Model?",
        [
            "Combines strengths of diverse algorithms",
            "Better generalization than individual models",
            "Reduces variance and overfitting",
            "More robust to data variations",
            "Provides probability estimates for uncertainty quantification"
        ]
    )
    
    # Slide 9: Model Performance
    add_metrics_slide(
        prs,
        "Model Performance (Test Accuracy)",
        [
            ("Logistic Regression", "89.11%"),
            ("Support Vector Machine", "90.99%"),
            ("Naive Bayes", "83.98%"),
            ("Random Forest", "93.20%"),
            ("Ensemble Model (Selected)", "91.36%")
        ]
    )
    
    # Slide 10: Results Analysis
    add_content_slide(
        prs,
        "Results Analysis",
        [
            "Ensemble model achieves 91.36% accuracy",
            "Strong separation between Low and Medium stress levels",
            "Some confusion between Medium and High stress (improvement area)",
            "High stress occasionally misclassified as Low (critical for intervention)",
            "Model provides probability distributions for each class"
        ]
    )
    
    # Slide 11: Feature Importance
    add_two_column_slide(
        prs,
        "Key Feature Insights",
        [
            "Sleep Hours - Highest impact (inverse correlation)",
            "Exam Pressure - Major stress contributor",
            "Financial Stress - Significant factor",
            "Assignment Deadlines - Time pressure increases stress"
        ],
        [
            "Exercise Hours - Protective factor (reduces stress)",
            "Social Activities - Negative correlation with stress",
            "Family Support - Reduces stress levels",
            "Workload Level - Strong predictor"
        ]
    )
    
    # Slide 12: System Architecture
    add_methodology_slide(
        prs,
        "System Architecture",
        [
            "Data Layer: Generation → Storage → Loading",
            "Preprocessing Layer: Cleaning → Balancing → Scaling",
            "Model Layer: Training → Ensemble → Persistence",
            "Evaluation Layer: Metrics → Visualizations → Reports",
            "Application Layer: Flask Web App → API → Recommendations"
        ]
    )
    
    # Slide 13: Web Application
    add_content_slide(
        prs,
        "Web Application Features",
        [
            "User-friendly input form for 10 features",
            "Real-time stress level prediction",
            "Probability distribution for each stress level",
            "Personalized motivational suggestions",
            "Responsive design for desktop and mobile",
            "RESTful API for programmatic access"
        ]
    )
    
    # Slide 14: Challenges & Solutions
    add_two_column_slide(
        prs,
        "Challenges & Solutions",
        [
            "Class Imbalance → SMOTE oversampling",
            "Feature Scaling → StandardScaler normalization",
            "Model Selection → Comprehensive comparison",
            "Overfitting → Ensemble approach"
        ],
        [
            "Generalization → Cross-validation ready",
            "Interpretability → Feature importance analysis",
            "Deployment → Flask web framework",
            "User Experience → Intuitive interface"
        ]
    )
    
    # Slide 15: Future Work
    add_content_slide(
        prs,
        "Future Enhancements",
        [
            "Validate on real-world student data",
            "Add temporal analysis for stress trends",
            "Integrate with Learning Management Systems",
            "Develop mobile application",
            "Implement deep learning models",
            "Add automated data collection",
            "Expand feature set (mental health history, support systems)"
        ]
    )
    
    # Slide 16: Conclusion
    add_conclusion_slide(
        prs,
        "Conclusion",
        [
            "Successfully developed end-to-end ML pipeline",
            "Ensemble model achieves 91.36% accuracy",
            "Comprehensive comparison of 4 ML algorithms",
            "User-friendly web application deployed",
            "Actionable insights and recommendations provided"
        ],
        "Ready for pilot deployment with real-world data"
    )
    
    prs.save(output_path)
    print(f"Professional presentation generated: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation("/home/sahu/Desktop/Student_S_detection/PROJECT_DECK.pptx")
